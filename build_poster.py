"""
Build poster.pptx — A2 scientific poster, fully editable in PowerPoint/LibreOffice.
Run: python build_poster.py
"""

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── COLOR PALETTE ───────────────────────────────────────────────────────────
BG       = RGBColor(0x08, 0x10, 0x1F)   # deep navy
BG2      = RGBColor(0x0D, 0x16, 0x28)
GOLD     = RGBColor(0xF5, 0xA6, 0x23)
GOLDLT   = RGBColor(0xFF, 0xD0, 0x6B)
RED      = RGBColor(0xE8, 0x45, 0x45)
TEAL     = RGBColor(0x00, 0xD4, 0xA8)
TEALLT   = RGBColor(0x4E, 0xFF, 0xD9)
WHITE    = RGBColor(0xEF, 0xF4, 0xFF)
DIM      = RGBColor(0x88, 0x92, 0xAA)
CARD     = RGBColor(0x10, 0x1A, 0x30)
DARK     = RGBColor(0x06, 0x0C, 0x18)

# ── SLIDE DIMENSIONS (A2 portrait) ──────────────────────────────────────────
W = Mm(420)
H = Mm(594)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)

# ── HELPERS ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        x, y, w, h
    )
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_name='Calibri', font_size=Pt(10),
                bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT,
                word_wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing.pt * 100)))
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox, tf


def add_multiline_textbox(slide, x, y, w, h, lines, word_wrap=True):
    """lines = list of dicts with keys: text, font, size, bold, italic, color, align, space_before"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = line.get('align', PP_ALIGN.LEFT)
        if line.get('space_before'):
            p.space_before = line['space_before']
        run = p.add_run()
        run.text = line.get('text', '')
        run.font.name  = line.get('font', 'Calibri')
        run.font.size  = line.get('size', Pt(9))
        run.font.bold  = line.get('bold', False)
        run.font.italic= line.get('italic', False)
        run.font.color.rgb = line.get('color', WHITE)
    return txBox, tf

# ════════════════════════════════════════════════════════════════════════════
# BACKGROUND
# ════════════════════════════════════════════════════════════════════════════
bg = add_rect(slide, 0, 0, W, H, fill_color=BG)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 1 · HEADER (0 → 110mm)
# ════════════════════════════════════════════════════════════════════════════
HDR_H = Mm(110)
hdr_bg = add_rect(slide, 0, 0, W, HDR_H, fill_color=BG2)

# Gold bottom border of header
add_rect(slide, 0, HDR_H - Mm(1), W, Mm(1.2), fill_color=GOLD)

# Eyebrow
add_textbox(slide,
    Mm(14), Mm(10), Mm(380), Mm(8),
    'WORLD DNA DAY · APRIL 25, 2026 · BRAC UNIVERSITY SOCIETY FOR BIOTECHNOLOGY (BUSB)',
    font_name='Calibri', font_size=Pt(8), bold=True, color=GOLD
)

# Main title
add_multiline_textbox(slide,
    Mm(14), Mm(20), Mm(260), Mm(75),
    [
        {'text': 'Rewriting Fate:', 'font': 'Calibri', 'size': Pt(34), 'bold': True, 'color': GOLD},
        {'text': 'How CRISPR-Cas9 Transforms Beta-Thalassemia', 'font': 'Calibri', 'size': Pt(34), 'bold': True, 'color': WHITE},
        {'text': 'from a Lifelong Burden into a', 'font': 'Calibri', 'size': Pt(34), 'bold': True, 'color': WHITE},
        {'text': 'One-Time Cure', 'font': 'Calibri', 'size': Pt(34), 'bold': True, 'color': TEAL},
    ]
)

# Tagline
add_multiline_textbox(slide,
    Mm(14), Mm(79), Mm(300), Mm(26),
    [
        {'text': 'Beta-thalassemia — a single HBB gene mutation — sentences millions to lifelong transfusions. In December 2023, the FDA approved Casgevy: the world\'s first CRISPR-Cas9 gene therapy. By permanently reactivating fetal hemoglobin, it freed 93% of trial patients from transfusion dependency. With Bangladesh carrying a ~7% carrier rate, this breakthrough is not abstract — it is urgent.',
         'font': 'Calibri', 'size': Pt(9.5), 'color': DIM},
    ]
)

# DNA helix label (right side decoration)
add_textbox(slide,
    Mm(295), Mm(38), Mm(110), Mm(12),
    'HBB GENE · CHROMOSOME 11p15.5',
    font_name='Calibri', font_size=Pt(8.5), bold=True, color=GOLD, align=PP_ALIGN.CENTER
)
# Decorative accent lines simulating helix
for i, (yoff, col) in enumerate([(0, GOLD), (8, TEAL), (16, GOLD), (24, TEAL)]):
    add_rect(slide, Mm(300 + i*0), Mm(20 + yoff), Mm(100), Mm(2),
             fill_color=col if i % 2 == 0 else TEAL)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 2 · STATS + DISEASE (110mm → 270mm)
# ════════════════════════════════════════════════════════════════════════════
S2_Y = HDR_H
S2_H = Mm(160)

# Stats panel background (left, 200mm wide)
add_rect(slide, 0, S2_Y, Mm(200), S2_H, fill_color=RGBColor(0x0C, 0x14, 0x24))
# Gold right border
add_rect(slide, Mm(200), S2_Y, Mm(1), S2_H, fill_color=RGBColor(0x60, 0x40, 0x10))

# STATS TITLE
add_textbox(slide, Mm(14), S2_Y + Mm(12), Mm(175), Mm(8),
    'BANGLADESH BURDEN & TRIAL OUTCOMES',
    font_name='Calibri', font_size=Pt(8), bold=True, color=GOLD)

# Stat 1 — 7%
add_textbox(slide, Mm(14), S2_Y + Mm(22), Mm(80), Mm(20),
    '7%', font_name='Calibri', font_size=Pt(40), bold=True, color=GOLD)
add_textbox(slide, Mm(14), S2_Y + Mm(42), Mm(175), Mm(10),
    'Carrier rate in Bangladesh — 1 in 14 people carry a mutated HBB allele.',
    font_name='Calibri', font_size=Pt(8), color=DIM)

# Divider
add_rect(slide, Mm(14), S2_Y + Mm(56), Mm(175), Mm(0.5), fill_color=RGBColor(0x40, 0x30, 0x08))

# Stat 2 — 14,000+
add_textbox(slide, Mm(14), S2_Y + Mm(60), Mm(120), Mm(16),
    '14,000+', font_name='Calibri', font_size=Pt(32), bold=True, color=GOLD)
add_textbox(slide, Mm(14), S2_Y + Mm(76), Mm(175), Mm(10),
    'New thalassemia major births annually in South Asia. Most in Bangladesh are undiagnosed.',
    font_name='Calibri', font_size=Pt(8), color=DIM)

add_rect(slide, Mm(14), S2_Y + Mm(90), Mm(175), Mm(0.5), fill_color=RGBColor(0x40, 0x30, 0x08))

# Stat 3 — 2-4x
add_textbox(slide, Mm(14), S2_Y + Mm(93), Mm(120), Mm(16),
    '2–4x / month', font_name='Calibri', font_size=Pt(28), bold=True, color=GOLD)
add_textbox(slide, Mm(14), S2_Y + Mm(109), Mm(175), Mm(10),
    'Blood transfusions required under traditional management — every month, for life.',
    font_name='Calibri', font_size=Pt(8), color=DIM)

add_rect(slide, Mm(14), S2_Y + Mm(123), Mm(175), Mm(0.5), fill_color=RGBColor(0x00, 0x50, 0x40))

# Stat 4 — 93% (teal)
add_textbox(slide, Mm(14), S2_Y + Mm(127), Mm(120), Mm(20),
    '93%', font_name='Calibri', font_size=Pt(38), bold=True, color=TEAL)
add_textbox(slide, Mm(14), S2_Y + Mm(147), Mm(175), Mm(10),
    'Patients transfusion-free after CRISPR therapy. CLIMB-Thal-111, NEJM 2023.',
    font_name='Calibri', font_size=Pt(8), color=TEALLT)

# ── DISEASE PANEL (right, 220mm wide) ────────────────────────────────────
DX = Mm(207)
DW = W - Mm(207) - Mm(10)

add_textbox(slide, DX, S2_Y + Mm(12), DW, Mm(8),
    'THE DISEASE',
    font_name='Calibri', font_size=Pt(8), bold=True, color=RED)

add_textbox(slide, DX, S2_Y + Mm(22), DW, Mm(16),
    'Beta-Thalassemia Major: One Gene, Catastrophic Cascade',
    font_name='Calibri', font_size=Pt(15), bold=True, color=WHITE)

add_textbox(slide, DX, S2_Y + Mm(40), DW, Mm(18),
    'Mutations in HBB (chromosome 11p15.5) cripple β-globin synthesis. Without functional hemoglobin, red blood cells collapse before leaving the bone marrow — ineffective erythropoiesis — causing severe hemolytic anemia and multi-organ failure.',
    font_name='Calibri', font_size=Pt(8.5), color=DIM)

# Thalassemic RBC box (left)
rbc_l = add_rect(slide, DX, S2_Y + Mm(62), Mm(95), Mm(50), fill_color=RGBColor(0x18, 0x08, 0x08))
rbc_l.line.color.rgb = RED
rbc_l.line.width = Pt(1)
add_textbox(slide, DX + Mm(3), S2_Y + Mm(64), Mm(89), Mm(7),
    'THALASSEMIC RBC', font_name='Calibri', font_size=Pt(7.5), bold=True, color=RED)
add_textbox(slide, DX + Mm(3), S2_Y + Mm(73), Mm(89), Mm(22),
    '◆ Irregular, fragile — destroyed in bone marrow\n◆ Severe hemolytic anemia\n◆ Ineffective erythropoiesis\n◆ Bone deformity, splenomegaly, cardiac failure',
    font_name='Calibri', font_size=Pt(7.5), color=RGBColor(0xFF, 0x88, 0x88))

# Post-CRISPR RBC box (right)
rbc_r = add_rect(slide, DX + Mm(101), S2_Y + Mm(62), Mm(95), Mm(50), fill_color=RGBColor(0x04, 0x18, 0x14))
rbc_r.line.color.rgb = TEAL
rbc_r.line.width = Pt(1)
add_textbox(slide, DX + Mm(104), S2_Y + Mm(64), Mm(89), Mm(7),
    'POST-CRISPR RBC · HbF ✓', font_name='Calibri', font_size=Pt(7.5), bold=True, color=TEAL)
add_textbox(slide, DX + Mm(104), S2_Y + Mm(73), Mm(89), Mm(22),
    '◆ Smooth biconcave discs with fetal hemoglobin\n◆ Full oxygen-carrying capacity restored\n◆ No transfusions required\n◆ Normal organ function maintained',
    font_name='Calibri', font_size=Pt(7.5), color=TEALLT)

# Symptom tags row
add_textbox(slide, DX, S2_Y + Mm(118), DW, Mm(8),
    'Untreated cascade: Severe Anemia  →  Bone Deformity  →  Splenomegaly  →  Cardiac Failure  →  Death in teens',
    font_name='Calibri', font_size=Pt(8), color=RED)

# Divider below section 2
add_rect(slide, 0, S2_Y + S2_H, W, Mm(0.6), fill_color=RGBColor(0x20, 0x20, 0x30))

# ════════════════════════════════════════════════════════════════════════════
# SECTION 3 · CRISPR MECHANISM HERO (270mm → 410mm)
# ════════════════════════════════════════════════════════════════════════════
S3_Y = S2_Y + S2_H
S3_H = Mm(140)

add_rect(slide, 0, S3_Y, W, S3_H, fill_color=RGBColor(0x0B, 0x15, 0x26))

add_textbox(slide, Mm(14), S3_Y + Mm(10), Mm(392), Mm(8),
    'CRISPR-Cas9 MECHANISM  ·  Casgevy (exagamglogene autotemcel)  ·  FDA + EMA APPROVED DECEMBER 8, 2023',
    font_name='Calibri', font_size=Pt(8), bold=True, color=TEAL)

add_textbox(slide, Mm(14), S3_Y + Mm(20), Mm(392), Mm(12),
    'Six Steps from Broken Gene to Permanent Cure',
    font_name='Calibri', font_size=Pt(18), bold=True, color=WHITE)

# ── 6 STEP BOXES ─────────────────────────────────────────────────────────
steps = [
    {
        'num': '01',
        'title': 'HSC Harvest',
        'body': 'Patient\'s own hematopoietic stem cells collected via apheresis. Plerixafor mobilization. No donor required.',
        'col': GOLD, 'bg': RGBColor(0x14, 0x10, 0x04)
    },
    {
        'num': '02',
        'title': 'Cas9 + gRNA Delivery',
        'body': 'Guide RNA directs Cas9 protein to the BCL11A erythroid enhancer (+58 region) — delivered ex vivo by electroporation.',
        'col': TEAL, 'bg': RGBColor(0x04, 0x14, 0x12)
    },
    {
        'num': '03',
        'title': 'BCL11A Enhancer Cut',
        'body': 'Cas9 introduces a double-strand break at the +58 enhancer. NHEJ repair creates indels, destroying the GATA1 binding motif.',
        'col': RED, 'bg': RGBColor(0x18, 0x06, 0x06)
    },
    {
        'num': '04',
        'title': 'HbF Genes Reactivate',
        'body': 'Without BCL11A in erythroid cells, HBG1/HBG2 (γ-globin) genes remain ON. Fetal hemoglobin (HbF) is produced permanently.',
        'col': TEAL, 'bg': RGBColor(0x04, 0x16, 0x12)
    },
    {
        'num': '05',
        'title': 'Conditioning + Reinfusion',
        'body': 'Myeloablative busulfan conditioning clears native marrow. Edited HSCs infused — they home to the bone marrow niche and engraft.',
        'col': GOLD, 'bg': RGBColor(0x14, 0x10, 0x04)
    },
    {
        'num': '06',
        'title': 'Transfusion-Free Life',
        'body': 'Every red blood cell the patient produces carries the edited genome. HbF compensates for absent HbA. One treatment. Permanent.',
        'col': TEAL, 'bg': RGBColor(0x04, 0x18, 0x14)
    },
]

BOX_W  = Mm(58)
BOX_H  = Mm(72)
GAP    = Mm(6)
STEP_Y = S3_Y + Mm(36)
STEP_X = Mm(14)

for i, step in enumerate(steps):
    bx = STEP_X + i * (BOX_W + GAP)
    box = add_rect(slide, bx, STEP_Y, BOX_W, BOX_H, fill_color=step['bg'])
    box.line.color.rgb = step['col']
    box.line.width = Pt(1.2)

    # Step number
    add_textbox(slide, bx + Mm(3), STEP_Y + Mm(3), BOX_W - Mm(6), Mm(7),
        step['num'], font_name='Calibri', font_size=Pt(7.5), bold=True, color=DIM)
    # Title
    add_textbox(slide, bx + Mm(3), STEP_Y + Mm(10), BOX_W - Mm(6), Mm(14),
        step['title'], font_name='Calibri', font_size=Pt(10), bold=True, color=step['col'])
    # Body
    add_textbox(slide, bx + Mm(3), STEP_Y + Mm(26), BOX_W - Mm(6), Mm(44),
        step['body'], font_name='Calibri', font_size=Pt(8), color=WHITE)

    # Arrow (except after last box)
    if i < 5:
        arr_x = bx + BOX_W + Mm(1)
        add_textbox(slide, arr_x, STEP_Y + BOX_H/2 - Mm(5), GAP - Mm(1), Mm(10),
            '→', font_name='Calibri', font_size=Pt(14), bold=True,
            color=TEAL if i == 4 else DIM, align=PP_ALIGN.CENTER)

# Insight callout
insight_y = STEP_Y + BOX_H + Mm(8)
ins = add_rect(slide, Mm(14), insight_y, W - Mm(28), Mm(18),
               fill_color=RGBColor(0x04, 0x14, 0x12))
ins.line.color.rgb = RGBColor(0x00, 0x80, 0x65)
ins.line.width = Pt(1)
add_textbox(slide, Mm(18), insight_y + Mm(2), W - Mm(36), Mm(14),
    'Key insight: Casgevy does NOT repair the mutated HBB gene. It disables BCL11A — the molecular switch that silences fetal hemoglobin after birth — permanently re-engaging a compensatory system already encoded in the patient\'s own genome. The edit is heritable in all descendent red blood cells through the engrafted stem cell population.',
    font_name='Calibri', font_size=Pt(8.5), color=WHITE)

# Divider
add_rect(slide, 0, S3_Y + S3_H, W, Mm(0.6), fill_color=RGBColor(0x20, 0x20, 0x30))

# ════════════════════════════════════════════════════════════════════════════
# SECTION 4 · CLINICAL EVIDENCE (410mm → 520mm)
# ════════════════════════════════════════════════════════════════════════════
S4_Y = S3_Y + S3_H
S4_H = Mm(110)
HALF = W // 2

# Left half
add_textbox(slide, Mm(14), S4_Y + Mm(10), HALF - Mm(20), Mm(8),
    'CLINICAL EVIDENCE', font_name='Calibri', font_size=Pt(8), bold=True, color=TEAL)
add_textbox(slide, Mm(14), S4_Y + Mm(20), HALF - Mm(20), Mm(14),
    'CLIMB-Thal-111 · Phase 3 Trial · NEJM December 2023',
    font_name='Calibri', font_size=Pt(13), bold=True, color=WHITE)

# Bar chart visual (using rectangles)
bars = [
    {'label': 'HbF BEFORE', 'pct': 0.12, 'val': '~8%',   'col': RED,  'bg': RGBColor(0x30, 0x08, 0x08)},
    {'label': 'HbF AFTER',  'pct': 0.65, 'val': '~43%',  'col': TEAL, 'bg': RGBColor(0x04, 0x20, 0x18)},
    {'label': 'TRANSFUSION-FREE', 'pct': 0.93, 'val': '93%', 'col': TEAL, 'bg': RGBColor(0x04, 0x24, 0x1A)},
]
BAR_MAX_H = Mm(55)
BAR_W     = Mm(35)
BAR_BASE  = S4_Y + Mm(95)
BAR_X     = Mm(14)
BAR_GAP   = Mm(15)

for i, bar in enumerate(bars):
    bx = BAR_X + i * (BAR_W + BAR_GAP)
    bar_h = int(BAR_MAX_H * bar['pct'])
    bar_y = BAR_BASE - bar_h
    # Background track
    add_rect(slide, bx, BAR_BASE - BAR_MAX_H, BAR_W, BAR_MAX_H, fill_color=RGBColor(0x10, 0x18, 0x28))
    # Fill
    b = add_rect(slide, bx, bar_y, BAR_W, bar_h, fill_color=bar['bg'])
    b.line.color.rgb = bar['col']
    b.line.width = Pt(1)
    # Value label
    add_textbox(slide, bx, bar_y - Mm(10), BAR_W, Mm(10),
        bar['val'], font_name='Calibri', font_size=Pt(13), bold=True,
        color=bar['col'], align=PP_ALIGN.CENTER)
    # Bottom label
    add_textbox(slide, bx, BAR_BASE + Mm(2), BAR_W, Mm(8),
        bar['label'], font_name='Calibri', font_size=Pt(6.5), bold=True,
        color=DIM, align=PP_ALIGN.CENTER)

add_textbox(slide, Mm(14), S4_Y + Mm(100), HALF - Mm(20), Mm(8),
    'n = 52 patients · Median follow-up 24.1 months · No off-target CRISPR editing detected · Effects durable at all timepoints',
    font_name='Calibri', font_size=Pt(7.5), color=DIM)

# Vertical divider
add_rect(slide, HALF, S4_Y, Mm(0.6), S4_H, fill_color=RGBColor(0x20, 0x20, 0x30))

# Right half — comparison
add_textbox(slide, HALF + Mm(10), S4_Y + Mm(10), HALF - Mm(20), Mm(8),
    'BEFORE vs. AFTER CRISPR', font_name='Calibri', font_size=Pt(8), bold=True, color=GOLD)
add_textbox(slide, HALF + Mm(10), S4_Y + Mm(20), HALF - Mm(20), Mm(12),
    'Traditional Treatment vs. One-Time CRISPR Cure',
    font_name='Calibri', font_size=Pt(12), bold=True, color=WHITE)

# Traditional box
tr = add_rect(slide, HALF + Mm(10), S4_Y + Mm(36), HALF//2 - Mm(15), Mm(68), fill_color=RGBColor(0x18, 0x06, 0x06))
tr.line.color.rgb = RED
tr.line.width = Pt(1)
add_textbox(slide, HALF + Mm(13), S4_Y + Mm(38), HALF//2 - Mm(20), Mm(7),
    'TRADITIONAL', font_name='Calibri', font_size=Pt(7.5), bold=True, color=RED)
add_textbox(slide, HALF + Mm(13), S4_Y + Mm(47), HALF//2 - Mm(20), Mm(55),
    '✕  Transfusion every 2–4 weeks, lifelong\n✕  Iron chelation — daily, toxic\n✕  BMT: only 30% find a matched donor\n✕  Graft-vs-host disease risk\n✕  Annual cost unaffordable in Bangladesh\n✕  Life expectancy severely reduced',
    font_name='Calibri', font_size=Pt(7.5), color=RGBColor(0xFF, 0x88, 0x88))

# CRISPR box
cx = HALF + Mm(10) + HALF//2 - Mm(10)
cr = add_rect(slide, cx, S4_Y + Mm(36), HALF//2 - Mm(15), Mm(68), fill_color=RGBColor(0x04, 0x18, 0x12))
cr.line.color.rgb = TEAL
cr.line.width = Pt(1)
add_textbox(slide, cx + Mm(3), S4_Y + Mm(38), HALF//2 - Mm(20), Mm(7),
    'CRISPR THERAPY', font_name='Calibri', font_size=Pt(7.5), bold=True, color=TEAL)
add_textbox(slide, cx + Mm(3), S4_Y + Mm(47), HALF//2 - Mm(20), Mm(55),
    '✓  One-time treatment — no maintenance\n✓  Autologous cells — no donor needed\n✓  No graft-vs-host disease risk\n✓  FDA + EMA approved Dec 8, 2023\n✓  93% transfusion-free, effects durable\n✓  Normal life expectancy projected',
    font_name='Calibri', font_size=Pt(7.5), color=TEALLT)

# Divider
add_rect(slide, 0, S4_Y + S4_H, W, Mm(0.6), fill_color=RGBColor(0x20, 0x20, 0x30))

# ════════════════════════════════════════════════════════════════════════════
# SECTION 5 · CASGEVY + BANGLADESH (520mm → 560mm)
# ════════════════════════════════════════════════════════════════════════════
S5_Y = S4_Y + S4_H
S5_H = Mm(55)

add_rect(slide, 0, S5_Y, HALF, S5_H, fill_color=RGBColor(0x0C, 0x12, 0x1E))
add_rect(slide, HALF, S5_Y, HALF, S5_H, fill_color=RGBColor(0x08, 0x12, 0x10))
add_rect(slide, HALF, S5_Y, Mm(0.6), S5_H, fill_color=RGBColor(0x40, 0x30, 0x08))

# Casgevy left
add_textbox(slide, Mm(14), S5_Y + Mm(8), Mm(100), Mm(7),
    "WORLD'S FIRST CRISPR MEDICINE", font_name='Calibri', font_size=Pt(8), bold=True, color=GOLD)

# FDA / EMA badges as colored boxes
fda = add_rect(slide, Mm(14), S5_Y + Mm(17), Mm(28), Mm(12), fill_color=RGBColor(0x0D, 0x22, 0x48))
fda.line.color.rgb = RGBColor(0x4A, 0x7B, 0xC8)
fda.line.width = Pt(1.5)
add_textbox(slide, Mm(14), S5_Y + Mm(19), Mm(28), Mm(9), 'FDA',
    font_name='Calibri', font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

ema = add_rect(slide, Mm(46), S5_Y + Mm(17), Mm(28), Mm(12), fill_color=RGBColor(0x0D, 0x2A, 0x1A))
ema.line.color.rgb = RGBColor(0x3A, 0x9A, 0x5C)
ema.line.width = Pt(1.5)
add_textbox(slide, Mm(46), S5_Y + Mm(19), Mm(28), Mm(9), 'EMA',
    font_name='Calibri', font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, Mm(80), S5_Y + Mm(17), Mm(60), Mm(12),
    'Approved\nDecember 8, 2023', font_name='Calibri', font_size=Pt(8.5), bold=False, color=GOLD)

add_textbox(slide, Mm(14), S5_Y + Mm(32), HALF - Mm(24), Mm(10),
    'Casgevy™  (exagamglogene autotemcel)  ·  Vertex Pharmaceuticals + CRISPR Therapeutics  ·  First CRISPR approval in human history  ·  Indication: TDT in patients ≥12 years',
    font_name='Calibri', font_size=Pt(7.5), color=DIM)

# Timeline bar
tl_y = S5_Y + Mm(44)
add_rect(slide, Mm(14), tl_y + Mm(2), HALF - Mm(28), Mm(2), fill_color=DIM)
for tx, label, col in [(Mm(14), '2012\nDiscovery', GOLD), (Mm(14) + (HALF - Mm(28))//2, '2020\nPhase 3', DIM), (Mm(14) + HALF - Mm(28), '2023\nAPPROVED', TEAL)]:
    add_rect(slide, tx - Mm(3), tl_y, Mm(6), Mm(6), fill_color=col)
    add_textbox(slide, tx - Mm(12), tl_y + Mm(7), Mm(24), Mm(10),
        label, font_name='Calibri', font_size=Pt(6.5), bold=True, color=col, align=PP_ALIGN.CENTER)

# Bangladesh right
add_textbox(slide, HALF + Mm(10), S5_Y + Mm(8), HALF - Mm(20), Mm(7),
    'PATH FORWARD FOR BANGLADESH', font_name='Calibri', font_size=Pt(8), bold=True, color=TEAL)
add_textbox(slide, HALF + Mm(10), S5_Y + Mm(17), HALF - Mm(20), Mm(36),
    '▶  Carrier screening: identify at-risk couples before pregnancy (~USD 15 per PCR test)\n'
    '▶  Prenatal diagnosis via CVS / amniocentesis with HBB genotyping\n'
    '▶  Genetic counseling at district hospital level — policy commitment required\n'
    '▶  Future: biosimilar gene therapies + regional manufacturing hubs in South Asia\n'
    '▶  Advocate for government-funded clinical programs for patients already born with TDT',
    font_name='Calibri', font_size=Pt(8), color=WHITE)

cta = add_rect(slide, HALF + Mm(10), S5_Y + Mm(46), HALF - Mm(20), Mm(7),
               fill_color=RGBColor(0x04, 0x18, 0x12))
cta.line.color.rgb = TEAL
cta.line.width = Pt(1.2)
add_textbox(slide, HALF + Mm(13), S5_Y + Mm(47), HALF - Mm(26), Mm(5),
    'The science is ready. Healthcare systems must catch up. Know your genes.',
    font_name='Calibri', font_size=Pt(8), bold=True, color=TEAL)

# ════════════════════════════════════════════════════════════════════════════
# FOOTER (560mm → 594mm)
# ════════════════════════════════════════════════════════════════════════════
FT_Y = S5_Y + S5_H
FT_H = H - FT_Y

add_rect(slide, 0, FT_Y, W, FT_H, fill_color=DARK)
add_rect(slide, 0, FT_Y, W, Mm(1.2), fill_color=GOLD)

add_textbox(slide, Mm(14), FT_Y + Mm(4), W - Mm(100), FT_H - Mm(6),
    'References:  1. Locatelli F et al. NEJM 2024;390:1663–1676  ·  2. Frangoul H et al. NEJM 2021;384:252–260  ·  3. FDA. Casgevy Approval. FDA.gov, Dec 8 2023  ·  '
    '4. Taher AT et al. Nat Rev Dis Primers 2021;7:24  ·  5. Doudna JA & Charpentier E. Science 2012;337:816–821  ·  6. Sankaran VG et al. Science 2008;322:1839–1842  ·  '
    '7. Uda M et al. PNAS 2008;105:1620–1625  ·  8. Weatherall DJ. Lancet 2010;391:155–167',
    font_name='Calibri', font_size=Pt(7), color=DIM)

# BUSB badge
busb = add_rect(slide, W - Mm(85), FT_Y + Mm(4), Mm(35), Mm(10), fill_color=GOLD)
add_textbox(slide, W - Mm(85), FT_Y + Mm(5), Mm(35), Mm(8),
    'BUSB', font_name='Calibri', font_size=Pt(9), bold=True,
    color=RGBColor(0x00, 0x00, 0x00), align=PP_ALIGN.CENTER)
add_textbox(slide, W - Mm(48), FT_Y + Mm(4), Mm(36), Mm(14),
    'World DNA Day 2026\nGenetic Disorder & Medication\nBRAC University',
    font_name='Calibri', font_size=Pt(7), color=DIM, align=PP_ALIGN.RIGHT)

# ── SAVE ────────────────────────────────────────────────────────────────────
out = 'F:/DNA Project/poster.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slide size: {prs.slide_width.mm:.0f}mm × {prs.slide_height.mm:.0f}mm (A2)')
print('Open in PowerPoint or LibreOffice Impress → File → Export as PDF')
